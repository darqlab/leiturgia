from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colours ───────────────────────────────────────────────────────────────────
C = {
    "bg":       RGBColor(0x0D, 0x1B, 0x3E),
    "bgCard":   RGBColor(0x1A, 0x2B, 0x5C),
    "bgDark":   RGBColor(0x09, 0x14, 0x28),
    "gold":     RGBColor(0xD4, 0xA8, 0x43),
    "goldLt":   RGBColor(0xF0, 0xC9, 0x6A),
    "white":    RGBColor(0xFF, 0xFF, 0xFF),
    "offWhite": RGBColor(0xE8, 0xEA, 0xF0),
    "muted":    RGBColor(0x8C, 0x9A, 0xB5),
}

FONT_H = "Georgia"
FONT_B = "Calibri"
W = Inches(10)
H = Inches(5.625)


def inches(n): return Inches(n)
def pt(n):     return Pt(n)


# ── Primitives ────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(1, inches(x), inches(y), inches(w), inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, font_size=20, font_face=FONT_B,
             color=None, bold=False, italic=False, align="center", valign="middle",
             word_wrap=False):
    if color is None:
        color = C["white"]
    txBox = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap

    from pptx.enum.text import MSO_ANCHOR
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.alignment = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)

    run = p.add_run()
    run.text = text
    run.font.size = pt(font_size)
    run.font.name = font_face
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return txBox


def bg(slide):
    add_rect(slide, 0, 0, 10, 5.625, C["bg"])


def bottom_bar(slide):
    add_rect(slide, 0, 5.55, 10, 0.075, C["gold"])


def top_bar(slide):
    add_rect(slide, 0, 0, 10, 0.75, C["bgDark"])
    add_rect(slide, 0, 0.75, 10, 0.05, C["gold"])


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ── Shared slide builders ─────────────────────────────────────────────────────
def add_title_slide(prs, section_name, church, date):
    slide = _blank_slide(prs)
    bg(slide)
    add_rect(slide, 0, 0, 10, 0.12, C["gold"])
    bottom_bar(slide)
    add_rect(slide, 0, 0.12, 0.08, 5.435, C["gold"])

    add_text(slide, church, 0.3, 0.25, 9, 0.45, font_size=14, color=C["muted"])
    add_text(slide, section_name, 0.3, 1.1, 9.4, 1.6,
             font_size=60, font_face=FONT_H, bold=True, align="left", color=C["white"])
    add_rect(slide, 0.3, 2.85, 6, 0.05, C["gold"])
    subtitle = f"Worship Program · {date}" if date else "Worship Program"
    add_text(slide, subtitle, 0.3, 3.0, 9, 0.5, font_size=20, color=C["goldLt"], italic=True)
    add_text(slide, '"Come, let us worship and bow down"  — Psalm 95:6',
             0.3, 3.7, 9, 0.5, font_size=15, color=C["muted"], italic=True)


def add_overview_slide(prs, section_label, time_str, items):
    """Flat item overview (used for Sabbath School)."""
    slide = _blank_slide(prs)
    bg(slide)
    top_bar(slide)
    bottom_bar(slide)

    add_text(slide, section_label.upper(), 0.3, 0.1, 6, 0.55, font_size=11, color=C["gold"], bold=True)
    add_text(slide, time_str, 0, 0.1, 9.7, 0.55, font_size=11, color=C["muted"], align="right")
    add_text(slide, "Order of Service", 0.4, 0.85, 9, 0.65,
             font_size=34, font_face=FONT_H, bold=True, color=C["white"], align="left")

    available = 3.875
    row_h  = min(0.55, available / max(len(items), 1))
    font_s = 15 if row_h >= 0.50 else 12

    for i, item in enumerate(items):
        y = 1.6 + i * row_h
        fill = C["bgCard"] if i % 2 == 0 else C["bg"]
        add_rect(slide, 0.4, y, 9.2, row_h, fill)
        add_rect(slide, 0.4, y, 0.06, row_h, C["gold"])
        item_type = item.get("type", "participant")
        label = item["title"]
        if item_type == "song":
            if item.get("hymn_number"):
                label += f"  #{item['hymn_number']}"
        else:
            detail = item.get("part") or item.get("subtitle", "")
            if detail and detail != item["title"]:
                label += f"  —  {detail}"
        add_text(slide, label, 0.55, y + row_h * 0.1, 9.0, row_h * 0.8,
                 font_size=font_s, color=C["offWhite"], align="left")


def add_ds_overview_slide(prs, ds):
    """Overview showing Divine Service subsections."""
    slide = _blank_slide(prs)
    bg(slide)
    top_bar(slide)
    bottom_bar(slide)

    add_text(slide, "DIVINE SERVICE", 0.3, 0.1, 6, 0.55, font_size=11, color=C["gold"], bold=True)
    add_text(slide, ds.get("time", ""), 0, 0.1, 9.7, 0.55, font_size=11, color=C["muted"], align="right")
    add_text(slide, "Order of Service", 0.4, 0.85, 9, 0.65,
             font_size=34, font_face=FONT_H, bold=True, color=C["white"], align="left")

    for i, sub in enumerate(ds.get("subsections", [])):
        y = 1.6 + i * 0.63
        fill = C["bgCard"] if i % 2 == 0 else C["bg"]
        add_rect(slide, 0.4, y, 9.2, 0.55, fill)
        add_rect(slide, 0.4, y, 0.06, 0.55, C["gold"])
        n = len(sub.get("items", []))
        label = f"{sub['title']}   ({n} items)"
        add_text(slide, label, 0.55, y + 0.07, 9.0, 0.4, font_size=15, color=C["offWhite"], align="left")


def add_section_divider_slide(prs, title):
    """Divider slide between Divine Service subsections."""
    slide = _blank_slide(prs)
    bg(slide)
    bottom_bar(slide)
    add_rect(slide, 0, 2.72, 10, 0.05, C["gold"])

    add_text(slide, title, 0, 1.6, 10, 1.3,
             font_size=48, font_face=FONT_H, bold=True, align="center", color=C["white"])
    add_rect(slide, 3.5, 2.95, 3, 0.04, C["gold"])


def add_item_slide(prs, item, index, total):
    slide = _blank_slide(prs)
    bg(slide)
    bottom_bar(slide)

    SLIDE_MID   = 2.8125
    TITLE_BOX   = 1.05
    TITLE_PAD   = 0.55
    TITLE_VIS   = TITLE_BOX - TITLE_PAD
    SUB_H       = 0.48
    SUB_PAD     = 0.08
    SEP_H       = 0.04
    STACK_ROW_H = 1.1
    STACK_GAP   = 0.2

    # Resolve fields from new typed schema (or legacy schema)
    item_type = item.get("type", "participant")

    # ── Content item: free-form text body, no participant row ─────────────────
    if item_type == "content":
        content_text = item.get("content", "")
        SLIDE_MID  = 2.8125
        TITLE_VIS  = 0.50
        SEP_H      = 0.04
        BODY_H     = 1.80
        total_vis  = TITLE_VIS + 0.30 + SEP_H + 0.22 + BODY_H
        title_y    = SLIDE_MID - total_vis / 2

        add_text(slide, item["title"], 0, title_y, 10, 1.05,
                 font_size=46, font_face=FONT_H, bold=True, align="center",
                 color=C["goldLt"], valign="top")
        sep_y = title_y + TITLE_VIS + 0.30
        add_rect(slide, 3.5, sep_y, 3, SEP_H, C["gold"])
        add_text(slide, content_text, 1.0, sep_y + SEP_H + 0.22, 8.0, BODY_H,
                 font_size=20, font_face=FONT_B, color=C["offWhite"],
                 align="center", valign="top", word_wrap=True)
        add_text(slide, f"{index + 1} / {total}", 8.5, 5.2, 1.3, 0.3,
                 font_size=10, color=C["muted"], align="right")
        return

    if item_type == "song":
        subtitle = f"Hymn #{item['hymn_number']}" if item.get("hymn_number") else ""
        p_name   = ""
        p_role   = ""
    elif "participants" in item and not item.get("type"):
        # Legacy schema migration for generator
        subtitle = item.get("subtitle", "")
        parts    = item.get("participants", [])
        p_name   = parts[0].get("name", "") if parts else ""
        p_role   = parts[0].get("role", "") if parts else ""
    else:
        # Participant item
        subtitle = item.get("part", "")
        p_name   = item.get("participant", "")
        p_role   = item.get("part", "")

    has_sub = bool(subtitle)
    p_count = 1 if p_name else 0

    GAP_TITLE_SUB = 0.18
    GAP_TITLE_SEP = 0.30
    GAP_SUB_SEP   = 0.20
    GAP_SEP_P     = 0.22

    subtitle_vis = (SUB_H - SUB_PAD) if has_sub else 0
    after_title  = (GAP_TITLE_SUB + subtitle_vis + GAP_SUB_SEP) if has_sub else GAP_TITLE_SEP
    p_block_h    = p_count * STACK_ROW_H + max(0, p_count - 1) * STACK_GAP
    total_vis    = TITLE_VIS + after_title + SEP_H + GAP_SEP_P + p_block_h
    title_box_y  = SLIDE_MID - total_vis / 2 - TITLE_PAD + 0.29

    add_text(slide, item["title"], 0, title_box_y, 10, TITLE_BOX,
             font_size=52, font_face=FONT_H, bold=True, align="center",
             color=C["white"], valign="top")

    cur_y = title_box_y + TITLE_PAD + TITLE_VIS

    if has_sub:
        cur_y += GAP_TITLE_SUB
        add_text(slide, subtitle, 1, cur_y - SUB_PAD, 8, SUB_H,
                 font_size=24, color=C["goldLt"], italic=True, align="center", valign="top")
        cur_y += subtitle_vis + GAP_SUB_SEP
    else:
        cur_y += GAP_TITLE_SEP

    add_rect(slide, 3.5, cur_y, 3, SEP_H, C["gold"])
    cur_y += SEP_H + GAP_SEP_P

    if p_name:
        add_text(slide, p_name, 0, cur_y, 10, 0.72,
                 font_size=36, font_face=FONT_B, bold=True, align="center",
                 color=C["offWhite"], valign="middle")
        if p_role:
            add_text(slide, p_role.upper(), 0, cur_y + 0.68, 10, 0.42,
                     font_size=20, color=C["gold"], align="center", valign="top")

    add_text(slide, f"{index + 1} / {total}", 8.5, 5.2, 1.3, 0.3,
             font_size=10, color=C["muted"], align="right")


def _hymn_slide_sequence(stanzas):
    """
    Build the slide display order.
    If a chorus/refrain exists, insert it after every verse:
      v1 → chorus → v2 → chorus → v3 → chorus …
    Bridges are placed in their original position and not repeated.
    Songs with no chorus are returned as-is.
    """
    repeat_types = {"chorus", "refrain"}
    verses   = [s for s in stanzas if s.get("type", "verse") not in repeat_types and s.get("type") != "bridge"]
    refrains = [s for s in stanzas if s.get("type", "verse") in repeat_types]

    if not refrains:
        return list(stanzas)

    chorus = refrains[0]
    seq = []
    for verse in verses:
        seq.append(verse)
        seq.append(chorus)
    return seq


def _stanza_label(stanza):
    t = stanza.get("type", "verse")
    n = stanza.get("number", "")
    if t == "chorus":
        return "CHORUS"
    if t == "refrain":
        return "REFRAIN"
    if t == "bridge":
        return f"BRIDGE {n}" if n else "BRIDGE"
    return f"VERSE {n}" if n else "VERSE"


def add_hymn_slides(prs, stanzas):
    sequence = _hymn_slide_sequence(stanzas)
    total    = len(sequence)

    for idx, stanza in enumerate(sequence):
        slide = _blank_slide(prs)
        bg(slide)
        bottom_bar(slide)

        label    = _stanza_label(stanza)
        is_chorus = stanza.get("type") in ("chorus", "refrain")
        label_color = C["gold"] if not is_chorus else C["goldLt"]

        # Label pill background + text — large enough for the operator to read at a glance
        add_rect(slide, 3.8, 0.14, 2.4, 0.46, C["bgDark"])
        add_rect(slide, 3.8, 0.14, 0.06, 0.46, label_color)
        add_text(slide, label, 3.86, 0.14, 2.34, 0.46,
                 font_size=16, color=label_color, bold=True, align="center")

        # "Next ›" hint in bottom-right corner so operator knows what's coming
        if idx < total - 1:
            next_label = _stanza_label(sequence[idx + 1])
            add_text(slide, f"Next › {next_label}", 7.2, 5.18, 2.6, 0.30,
                     font_size=9, color=C["muted"], align="right")

        # Dynamic sizing: scale down as line count grows so all lines fit
        CONTENT_TOP = 0.68   # below label pill
        CONTENT_BOT = 5.28   # above dots
        AVAILABLE   = CONTENT_BOT - CONTENT_TOP   # 4.60"

        lines = stanza["lines"]
        n     = len(lines)

        # (line_h, gap, font_size) by line count
        if n <= 4:
            LINE_H, LINE_GAP, FONT_S = 0.80, 0.20, 26
        elif n == 5:
            LINE_H, LINE_GAP, FONT_S = 0.68, 0.14, 22
        elif n == 6:
            LINE_H, LINE_GAP, FONT_S = 0.58, 0.10, 19
        elif n == 7:
            LINE_H, LINE_GAP, FONT_S = 0.50, 0.08, 17
        else:
            LINE_H, LINE_GAP, FONT_S = 0.44, 0.06, 15

        block_h = LINE_H * n + LINE_GAP * (n - 1)
        start_y = CONTENT_TOP + (AVAILABLE - block_h) / 2

        for i, line in enumerate(lines):
            y = start_y + i * (LINE_H + LINE_GAP)
            add_text(slide, line, 0.3, y, 9.4, LINE_H,
                     font_size=FONT_S, font_face=FONT_H, align="center",
                     color=C["white"], valign="middle")

        # Dot progress bar — chorus dots are gold-tinted so operator can see the pattern
        dot_w   = 0.14
        dot_gap = 0.09
        total_w = dot_w * total + dot_gap * (total - 1)
        start_x = (10 - total_w) / 2
        for di in range(total):
            x = start_x + di * (dot_w + dot_gap)
            is_chorus_dot = sequence[di].get("type") in ("chorus", "refrain")
            if di == idx:
                dot_color = C["gold"]
            elif is_chorus_dot:
                dot_color = C["goldLt"]
            else:
                dot_color = C["muted"]
            shape = slide.shapes.add_shape(9, Inches(x), Inches(5.35), Inches(dot_w), Inches(dot_w))
            shape.fill.solid()
            shape.fill.fore_color.rgb = dot_color
            shape.line.fill.background()



def add_service_team_slide(prs, service_team, church):
    slide = _blank_slide(prs)
    bg(slide)
    top_bar(slide)
    bottom_bar(slide)

    add_text(slide, "DIVINE SERVICE", 0.3, 0.1, 6, 0.55, font_size=11, color=C["gold"], bold=True)
    add_text(slide, church, 0, 0.1, 9.7, 0.55, font_size=11, color=C["muted"], align="right")
    add_text(slide, "Service Team", 0.4, 0.85, 9, 0.65,
             font_size=34, font_face=FONT_H, bold=True, color=C["white"], align="left")

    row_h = min(0.62, 3.8 / max(len(service_team), 1))
    for i, member in enumerate(service_team):
        y = 1.6 + i * row_h
        fill = C["bgCard"] if i % 2 == 0 else C["bg"]
        add_rect(slide, 0.4, y, 9.2, row_h - 0.04, fill)
        add_rect(slide, 0.4, y, 0.06, row_h - 0.04, C["gold"])
        add_text(slide, member.get("role", "").upper(), 0.6, y + 0.04, 3.2, row_h * 0.38,
                 font_size=10, color=C["gold"], bold=True, align="left")
        add_text(slide, member.get("name", ""), 0.6, y + row_h * 0.42, 8.5, row_h * 0.52,
                 font_size=16, color=C["offWhite"], align="left")


def add_closing_slide(prs, section_name, church, next_service=None):
    slide = _blank_slide(prs)
    bg(slide)
    add_rect(slide, 0, 0, 10, 0.12, C["gold"])
    add_rect(slide, 0, 5.505, 10, 0.12, C["gold"])
    add_rect(slide, 0, 0.12, 0.08, 5.385, C["gold"])
    add_rect(slide, 9.92, 0.12, 0.08, 5.385, C["gold"])

    add_text(slide, section_name, 0.4, 1.3, 9.2, 1.1,
             font_size=56, font_face=FONT_H, bold=True, align="center", color=C["white"])
    add_rect(slide, 3.5, 2.55, 3, 0.05, C["gold"])
    ending = next_service if next_service else "Ends Here"
    add_text(slide, ending, 0.4, 2.75, 9.2, 0.55,
             font_size=20, color=C["gold"], italic=True, align="center")
    add_text(slide, f"{church}  ·  Sabbath Worship", 0.4, 5.0, 9.2, 0.35,
             font_size=13, color=C["muted"], align="center")


# ── Entry point ───────────────────────────────────────────────────────────────
def generate_pptx(program: dict, section: dict, output_path: str):
    """Generate a PPTX for any service program.

    Args:
        program: full program dict (church, date, service_team, …)
        section: one entry from program["service_programs"]
        output_path: where to save the .pptx
    """
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    church = program["church"]
    date   = program.get("date", "")
    name   = section.get("name", "")
    time   = section.get("time", "")
    items  = section.get("items", [])

    # Determine closing message based on next program in list
    programs = program.get("service_programs", [])
    idx      = next((i for i, sp in enumerate(programs) if sp["id"] == section["id"]), -1)
    next_sp  = programs[idx + 1] if 0 <= idx < len(programs) - 1 else None
    closing_next = f"Ends Here · {next_sp['name']} Follows" if next_sp else None

    add_title_slide(prs, name, church, date)
    add_overview_slide(prs, name, time, items)

    for i, item in enumerate(items):
        add_item_slide(prs, item, i, len(items))
        if item.get("lyrics"):
            add_hymn_slides(prs, item["lyrics"])

    if program.get("service_team"):
        add_service_team_slide(prs, program["service_team"], church)

    add_closing_slide(prs, name, church, closing_next)

    prs.save(output_path)
    print(f"✅ Generated: {output_path}")
